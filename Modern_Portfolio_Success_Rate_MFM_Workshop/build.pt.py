from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# --------- USER SETTINGS ----------
OUT_PPTX = "Retirement_Success_Project.pptx"
WR_PNG = "wr_frontier.png"

# Fill these with your final numbers (edit if needed)
table_4pct = [
    ("Baseline", "89.61%"),
    ("Put", "89.97%"),
    ("Collar", "72.41%"),
    ("Guardrail", "94.25%"),
    ("Guardrail + Put", "95.29%"),
    ("Best Grid (Guardrail+Put)", "98.11%"),
]

table_wr_thresholds = [
    ("≥90%", "3.75%", "4.50%", "5.25%"),
    ("≥95%", "3.00%", "3.75%", "4.50%"),
]
# ----------------------------------

def add_title_and_bullets(slide, title, bullets):
    title_box = slide.shapes.title
    title_box.text = title

    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    for i, b in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(20)

def add_table(slide, left, top, width, height, data, col_names):
    rows = len(data) + 1
    cols = len(col_names)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # header
    for j, name in enumerate(col_names):
        cell = table.cell(0, j)
        cell.text = name
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # body
    for i, row in enumerate(data, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            if j > 0:
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    return table

def main():
    prs = Presentation()

    # Slide 1: Goal & Setup
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_bullets(
        slide,
        "Retirement Portfolio Success Rate Optimization (50-year horizon)",
        [
            "Goal: maximize probability of not running out of money over 50 years",
            "Initial wealth: $1,000,000",
            "Baseline spending: 4% of initial wealth, inflated by 3% yearly",
            "Monte Carlo: multivariate normal returns using means/vols/correlations from CSV inputs",
            "Success = portfolio value never falls below $0",
        ],
    )

    # Slide 2: Baseline allocation search
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_bullets(
        slide,
        "Step 1 — Search 5,000 Allocations (Long-only)",
        [
            "Sampled 5,000 random portfolios (Dirichlet weights)",
            "Simulated 10,000 retirement paths per portfolio",
            "Selected portfolio with highest baseline success rate",
            "Compared top 10% vs bottom 10% allocations to infer what drives success",
        ],
    )

    # Slide 3: Strategy tests @ 4% WR + table
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    slide.shapes.title.text = "Steps 2–6 — Strategy Tests & Tuning (WR=4%)"

    # table left
    add_table(
        slide,
        left=Inches(0.7),
        top=Inches(1.6),
        width=Inches(5.5),
        height=Inches(3.8),
        data=table_4pct,
        col_names=["Strategy", "Success @ 4% WR"],
    )

    # notes right
    tx = slide.shapes.add_textbox(Inches(6.4), Inches(1.6), Inches(3.0), Inches(3.8)).text_frame
    tx.word_wrap = True
    p = tx.paragraphs[0]
    p.text = "What we tested:"
    p.font.bold = True
    p.font.size = Pt(18)
    for line in [
        "• Protective Put (early years)",
        "• Collar (early years)",
        "• Guardrail spending: skip inflation after negative return (per-path)",
        "• Guardrail + Put",
        "• Grid search: hedge years / strike / premium",
    ]:
        q = tx.add_paragraph()
        q.text = line
        q.font.size = Pt(16)

    # Slide 4: WR frontier graph + threshold table
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    slide.shapes.title.text = "Step 7 — Success vs Withdrawal Rate (WR Frontier)"

    # image left
    if Path(WR_PNG).exists():
        slide.shapes.add_picture(WR_PNG, Inches(0.6), Inches(1.4), width=Inches(6.2))
    else:
        box = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(6.2), Inches(4.5)).text_frame
        box.text = f"Missing image: {WR_PNG}\nRun make_wr_frontier_plot.py first."
        box.paragraphs[0].font.size = Pt(18)

    # table right
    add_table(
        slide,
        left=Inches(7.0),
        top=Inches(1.6),
        width=Inches(2.8),
        height=Inches(2.0),
        data=table_wr_thresholds,
        col_names=["Target", "Baseline", "Guardrail", "Guardrail+Put"],
    )

    # takeaway
    takeaway = slide.shapes.add_textbox(Inches(7.0), Inches(3.8), Inches(2.8), Inches(2.0)).text_frame
    takeaway.word_wrap = True
    takeaway.paragraphs[0].text = "Takeaway"
    takeaway.paragraphs[0].font.bold = True
    takeaway.paragraphs[0].font.size = Pt(18)
    t1 = takeaway.add_paragraph()
    t1.text = "Guardrail shifts the WR frontier upward; adding a put shifts it further."
    t1.font.size = Pt(14)

    # Slide 5: Robustness
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_bullets(
        slide,
        "Step 8 — Robustness / Validation",
        [
            "Sanity A: multi-seed (10k paths) → strategy ranking stable",
            "Sanity B: 50k paths → thresholds remain similar (reduced Monte Carlo noise)",
            "Determinism check: same seed reproduces identical thresholds",
            "Stress tests: lower mean returns (−1%) and higher volatility (×1.25)",
        ],
    )

    prs.save(OUT_PPTX)
    print("Saved:", OUT_PPTX)

if __name__ == "__main__":
    main()